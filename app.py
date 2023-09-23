import streamlit as st
import docx
import PyPDF2
from pptx import Presentation
import re
import os
import openai

def query_openai(api_key, messages):
    openai.api_key = api_key
    response = openai.ChatCompletion.create(model="gpt-3.5-turbo", messages=messages)
    return response.choices[0].message["content"]

def extract_text_from_txt(file):
    try:
        return file.getvalue().decode('utf-8')
    except Exception as e:
        st.error(f"Error processing text file. Error: {e}")
        return None

def extract_text_from_docx(file):
    try:
        doc = docx.Document(file)
        return ' '.join([para.text for para in doc.paragraphs])
    except Exception as e:
        st.error(f"Error processing docx file. Error: {e}")
        return None

def extract_text_from_pdf(file):
    try:
        pdf_reader = PyPDF2.PdfFileReader(file)
        text = ""
        for page_num in range(pdf_reader.numPages):
            text += pdf_reader.getPage(page_num).extractText()
        return text
    except Exception as e:
        st.error(f"Error processing pdf file. Error: {e}")
        return None

def extract_text_from_ppt(file):
    try:
        prs = Presentation(file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text
        return text
    except Exception as e:
        st.error(f"Error processing ppt file. Error: {e}")
        return None

def extract_text(uploaded_file):
    try:
        if uploaded_file.type == "text/plain":
            return extract_text_from_txt(uploaded_file)
        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            return extract_text_from_docx(uploaded_file)
        elif uploaded_file.type == "application/pdf":
            return extract_text_from_pdf(uploaded_file)
        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return extract_text_from_ppt(uploaded_file)
        else:
            return None
    except Exception as e:
        st.error(f"Error processing {uploaded_file.name}. Error: {e}")
        return None

def find_keyword_in_text(keyword, text):
    matches = re.finditer(keyword, text, re.IGNORECASE)
    snippets = []
    for match in matches:
        start_index = max(0, match.start() - SNIPPET_LENGTH)
        end_index = min(len(text), match.end() + SNIPPET_LENGTH)
        snippet = text[start_index:end_index]
        snippets.append(snippet)
    return snippets

def extract_insights(text):
    insights = ""
    segments = text.split('. ')  # Split by sentences
    
    # Join sentences until they approach the segment size
    i = 0
    while i < len(segments):
        segment = segments[i]
        while len(segment) < SEGMENT_SIZE and i < len(segments) - 1:
            i += 1
            segment += ". " + segments[i]
        
        try:
            response = openai.Completion.create(
                model="text-davinci-003",
                prompt=f"Provide insights on the following transcript segment: {segment}",
                max_tokens=MAX_TOKENS
            )
            insights += response.choices[0].text.strip() + " "
        except Exception as e:
            st.error(f"OpenAI API error: {e}")
        
        i += 1
    
    return insights.strip()

# Constants
MAX_TOKENS = 200
SEGMENT_SIZE = 1000
SNIPPET_LENGTH = 50

st.title("Transcript Analysis Tool")

api_key = st.text_input("API Key", type="password")  # Moved API Key input here

uploaded_files = st.file_uploader(
    "Choose transcript files (.txt, .docx, .pdf, .ppt)",
    type=["txt", "docx", "pdf", "ppt"],
    accept_multiple_files=True,
)

# New file size check and filtering
accepted_files = [file for file in uploaded_files if file.size <= 10e6]  # 10 MB limit
rejected_files = [file for file in uploaded_files if file.size > 10e6]

for rejected_file in rejected_files:
    st.error(f"{rejected_file.name} is too large. Please upload smaller files.")

guiding_questions = st.text_area("Enter the guiding questions or keywords (separated by commas)")

# Dictionary to store text content for each accepted file
file_contents = {}

if accepted_files:
    with st.expander("Uploaded Files & Previews"):  # Updated to use expander
        for accepted_file in accepted_files:
            text_content = extract_text(accepted_file)
            file_contents[accepted_file.name] = text_content
            st.write(f"Contents of {accepted_file.name}: {text_content[:500]}...")

if guiding_questions:
    with st.expander("Keyword Matches"):
        keywords = [keyword.strip() for keyword in guiding_questions.split(",")]

        for keyword in keywords:
            for file_name, text_content in file_contents.items():
                snippets = find_keyword_in_text(keyword, text_content)
                if snippets:
                    st.write(f"Found {len(snippets)} instances of '{keyword}' in {file_name}:")
                    for snippet in snippets:
                        st.write(f"...{snippet}...")

# Display insights
with st.expander("Extracted Insights"):
    for file_name, text_content in file_contents.items():
        insights = extract_insights(text_content)
        st.write(f"Insights from {file_name}:")
        st.write(insights)
